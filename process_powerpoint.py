#!/usr/bin/env python3
"""
PowerPoint Anonymization Module
Handles .pptx and .ppt files for the DOCX Anonymizer app
"""

from pptx import Presentation
from pathlib import Path
import logging
import re


def strip_pptx_metadata(prs):
    """
    Strip ALL metadata from PowerPoint file.

    Similar to Word metadata stripping - removes:
    - Author
    - Title
    - Subject
    - Keywords
    - Comments
    - Company
    """
    props = prs.core_properties

    props.author = ""
    props.last_modified_by = ""
    props.title = ""
    props.subject = ""
    props.keywords = ""
    props.comments = ""
    props.category = ""
    props.content_status = ""
    props.identifier = ""

    # Clear company (important for SEC filings)
    if hasattr(props, 'company'):
        props.company = ""

    props.revision = 1
    if hasattr(props, 'version'):
        props.version = None

    return prs


def remove_all_images_pptx(prs):
    """
    Remove ALL images from PowerPoint slides.

    Returns count of removed images.
    """
    removed_count = 0

    for slide in prs.slides:
        # Find all picture shapes
        shapes_to_remove = []
        for shape in slide.shapes:
            # Check if shape is a picture
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                shapes_to_remove.append(shape)

        # Remove pictures (must be done after iteration)
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)
            removed_count += 1

    return removed_count


def anonymize_text_pptx(text, alias_map, sorted_keys, compiled_patterns, track_details=False):
    """
    Apply anonymization replacements with case matching using SINGLE-PASS regex (v2.1).

    Reuses same logic and pattern format as Word processor.

    Returns:
        If track_details=False: (text, replacements)
        If track_details=True: (text, replacements, details_dict)
    """
    replacements = 0

    # Extract combined pattern and lookup map
    combined_pattern = compiled_patterns.get('combined')
    lookup = compiled_patterns.get('lookup')

    # BACKWARD COMPATIBILITY: Handle old compiled_patterns format
    if combined_pattern is None or lookup is None:
        # Old format - use legacy multi-pass algorithm
        result = anonymize_text_pptx_legacy(text, alias_map, sorted_keys, compiled_patterns)
        if track_details:
            return result[0], result[1], {}
        return result

    # Track which originals were replaced (v2.1 feature)
    details = {} if track_details else None

    # SINGLE-PASS REPLACEMENT (v2.1 performance optimization)
    def replace_match(match):
        nonlocal replacements
        matched_text = match.group(0)

        # Look up the replacement using lowercase match
        matched_lower = matched_text.lower()
        if matched_lower not in lookup:
            return matched_text  # Safe fallback

        original, replacement = lookup[matched_lower]

        # Track this replacement (v2.1)
        if track_details:
            details[original] = details.get(original, 0) + 1

        # Preserve case pattern
        if matched_text.isupper():
            replacements += 1
            return replacement.upper()
        elif matched_text.islower():
            replacements += 1
            return replacement.lower()
        elif matched_text[0].isupper():
            replacements += 1
            return replacement.capitalize()
        else:
            replacements += 1
            return replacement

    # Single regex pass replaces ALL patterns at once
    text = combined_pattern.sub(replace_match, text)

    if track_details:
        return text, replacements, details
    return text, replacements


def merge_details(details1, details2):
    """
    Merge two replacement details dictionaries (v2.1 helper).
    """
    if details1 is None:
        return details2 if details2 else {}
    if details2 is None:
        return details1

    merged = details1.copy()
    for original, count in details2.items():
        merged[original] = merged.get(original, 0) + count
    return merged


def anonymize_text_pptx_legacy(text, alias_map, sorted_keys, compiled_patterns):
    """
    Legacy multi-pass anonymization (kept for backward compatibility).
    """
    replacements = 0

    for original in sorted_keys:
        replacement = alias_map[original]

        def replace_with_case(match):
            nonlocal replacements
            matched_text = match.group(0)

            # Preserve case pattern
            if matched_text.isupper():
                replacements += 1
                return replacement.upper()
            elif matched_text.islower():
                replacements += 1
                return replacement.lower()
            elif matched_text[0].isupper():
                replacements += 1
                return replacement.capitalize()
            else:
                replacements += 1
                return replacement

        pattern = compiled_patterns[original]
        text = pattern.sub(replace_with_case, text)

    return text, replacements


def anonymize_pptx(pptx_path, alias_map, sorted_keys, compiled_patterns, track_details=False):
    """
    Anonymize all text in PowerPoint file (v2.1 with optional tracking).

    Processes:
    - Slide text frames
    - Tables
    - Notes
    - Shapes

    Returns:
        If track_details=False: (prs, total_replacements)
        If track_details=True: (prs, total_replacements, details_dict)
    """
    prs = Presentation(pptx_path)
    total_replacements = 0
    document_details = {} if track_details else None

    # Create tracking wrapper
    def anonymize_with_tracking(text, alias_map, sorted_keys, compiled_patterns):
        nonlocal document_details
        if track_details:
            new_text, count, details = anonymize_text_pptx(text, alias_map, sorted_keys, compiled_patterns, track_details=True)
            document_details = merge_details(document_details, details)
            return new_text, count
        else:
            return anonymize_text_pptx(text, alias_map, sorted_keys, compiled_patterns)

    for slide in prs.slides:
        # Process all shapes with text
        for shape in slide.shapes:
            # Text frames (title, text boxes, etc.)
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            new_text, count = anonymize_with_tracking(
                                run.text, alias_map, sorted_keys, compiled_patterns
                            )
                            if count > 0:
                                run.text = new_text
                                total_replacements += count

            # Tables
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text:
                                    new_text, count = anonymize_with_tracking(
                                        run.text, alias_map, sorted_keys, compiled_patterns
                                    )
                                    if count > 0:
                                        run.text = new_text
                                        total_replacements += count

        # Process notes (speaker notes)
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, 'notes_text_frame'):
                for paragraph in notes_slide.notes_text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            new_text, count = anonymize_with_tracking(
                                run.text, alias_map, sorted_keys, compiled_patterns
                            )
                            if count > 0:
                                run.text = new_text
                                total_replacements += count

    if track_details:
        return prs, total_replacements, document_details
    return prs, total_replacements


def process_single_pptx(input_path, output_path, alias_map, sorted_keys, compiled_patterns, logger, remove_images=True, track_details=False):
    """
    Process a single PowerPoint file: anonymize + strip metadata + optional image removal.

    Args:
        input_path: Path to input .pptx file (string or Path object)
        output_path: Path for output .pptx file (string or Path object)
        alias_map: Dictionary of original → replacement mappings
        sorted_keys: Sorted list of alias_map keys
        compiled_patterns: Pre-compiled regex patterns
        logger: Logger instance
        remove_images: If True, removes all images from presentation
        track_details: If True, return detailed replacement tracking (v2.1)

    Returns:
        If track_details=False: (replacements, images_removed)
        If track_details=True: (replacements, images_removed, details_dict)
    """
    # Convert to Path objects if strings (for backward compatibility)
    from pathlib import Path
    input_path = Path(input_path) if isinstance(input_path, str) else input_path
    output_path = Path(output_path) if isinstance(output_path, str) else output_path

    logger.info(f"Processing: {input_path.name}")

    try:
        # Load and anonymize PowerPoint with optional tracking
        if track_details:
            prs, replacements, details = anonymize_pptx(input_path, alias_map, sorted_keys, compiled_patterns, track_details=True)
        else:
            prs, replacements = anonymize_pptx(input_path, alias_map, sorted_keys, compiled_patterns)

        # Remove all images (if requested)
        images_removed = 0
        if remove_images:
            images_removed = remove_all_images_pptx(prs)

        # Strip ALL metadata (CRITICAL)
        prs = strip_pptx_metadata(prs)

        # Save
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

        logger.info(f"  ✓ {replacements} replacements, {images_removed} images removed")

        if track_details:
            return replacements, images_removed, details
        return replacements, images_removed

    except Exception as e:
        logger.error(f"  ❌ Error: {e}")
        if track_details:
            return 0, 0, {}
        return 0, 0
